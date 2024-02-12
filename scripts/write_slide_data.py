import json
import re
import sys
import os
from pptx import Presentation

if len(sys.argv) < 3:
    raise Exception('Required arguments are: path to pptx and path to a folder with slide images')

pptx_path = sys.argv[1]
images_path = sys.argv[2]

prs = Presentation(pptx_path)

prs_title = prs.core_properties.title

if prs_title is None or len(prs_title) == 0 or prs_title == "PowerPoint Presentation":
    file_name = os.path.basename(pptx_path)
    prs_title = file_name.split('.')[0]

slide_info = []
max_length = 100

for index, slide in enumerate(prs.slides):
    title = None
    title_shape = slide.shapes.title
    # If there is a dedicated title then take it, otherwise find a first shape with text and use first part of that text
    if title_shape and hasattr(title_shape, "text") and len(title_shape.text) > 0:
        title = re.sub(r"\s+", " ", title_shape.text).strip()
    else:
        for shape in slide.shapes:
            if hasattr(shape, "text") and len(shape.text.strip()) > 0:
                title = re.sub("[^0-9a-zA-Z ]+", "", shape.text.strip())
                title = re.sub(r"\s+", " ", title).strip()
                break

    if title is None:
        title = "Untitled " + str(index)

    if len(title) > max_length:
        # TODO: may be better to find the closest space
        title = title[0:max_length]

    # Get slide notes if defined
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    slide_info.append({"title": title, "notes": notes})

folder_content = os.listdir(images_path)

slides_data = []

for entry in folder_content:
    if os.path.isdir(os.path.join(images_path, entry)):
        continue

    (_, ext) = os.path.splitext(entry)
    # Names are supposed to be like Slide-0.jpg, Slide-1.jpg etc.
    index_match = re.findall(r"-(\d+)" + re.escape(ext) + r"$", entry)
    if len(index_match) == 0:
        # When there is a single slide there is no index in title
        index = 0
    else:
        index = int(index_match[0])

    slide = slide_info[index]

    slides_data.append({
        "index": index,
        "title": slide["title"],
        "notes": slide["notes"],
        "image": entry
    })


def sort_by_index(slide):
    return slide["index"]


slides_data.sort(key=sort_by_index)
data_file = None
try:
    data_file = open(os.path.join(images_path, 'data.json'), "w")
    data_file.write(json.dumps({"title": prs_title, "slides": slides_data}, ensure_ascii=False))
finally:
    if data_file is not None:
        data_file.close()
